"""slide_generator — Generate PPTX presentations from experiment results."""

from dataclasses import dataclass, field
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


@dataclass
class SlideContent:
    title: str
    bullet_points: list[str] = field(default_factory=list)
    figure_path: str | None = None
    notes: str = ""


@dataclass
class SlideSpec:
    title: str
    subtitle: str
    author: str
    date: str
    slides: list[SlideContent]


def generate_deck(spec: SlideSpec, output_path: str = "presentation.pptx") -> Presentation:
    """Build PPTX from SlideSpec."""
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = spec.title
    slide.placeholders[1].text = f"{spec.subtitle}\n{spec.author} — {spec.date}"

    for sc in spec.slides:
        if sc.figure_path and Path(sc.figure_path).exists():
            # Figure slide (blank layout)
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # title only
            slide.shapes.title.text = sc.title
            slide.shapes.add_picture(sc.figure_path, Inches(1), Inches(1.8), Inches(8), Inches(5))
        else:
            # Bullet slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
            slide.shapes.title.text = sc.title
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(sc.bullet_points):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

        if sc.notes:
            slide.notes_slide.notes_text_frame.text = sc.notes

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return prs


if __name__ == "__main__":
    spec = SlideSpec(
        title="HYP-001: IG Spread Momentum",
        subtitle="Credit Research",
        author="AgentDS",
        date="2025-03-03",
        slides=[
            SlideContent(title="Hypothesis", bullet_points=[
                "H₀: Spread momentum has no predictive power",
                "H₁: 5-day spread momentum predicts excess returns",
                "Method: GARCH + vectorized backtest",
            ]),
            SlideContent(title="Results", bullet_points=[
                "OOS Sharpe: 0.82",
                "t-stat: 1.91 (borderline significant)",
                "Signal works in low-vol regimes only",
                "Transaction costs survive at 5bps",
            ]),
            SlideContent(title="Cumulative Returns", figure_path="figures/drawdown.png"),
        ],
    )
    generate_deck(spec, "../../artifacts/reports/hyp_001_slides.pptx")
    print("Deck generated.")
